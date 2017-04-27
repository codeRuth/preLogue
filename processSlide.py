from pptx import Presentation
from pptx.enum.shapes import PP_PLACEHOLDER
from pptx.exc import PackageNotFoundError

from libkey.keywordr import keywordr as k

import sys
reload(sys)
sys.setdefaultencoding('utf8')

class ProcessSlide(object):
    def __init__(self, inputPresentation):
        self.inputPresentation = inputPresentation

    def getTitle(self):
        titleObj = list()
        listObj = dict()
        try:
            p = Presentation(self.inputPresentation)
            for index, slide in enumerate(p.slides):
                for indexSlide, shape in enumerate(slide.shapes):
                    if shape.is_placeholder:
                        ph = shape.placeholder_format
                        if ph.type == PP_PLACEHOLDER.TITLE or ph.type == PP_PLACEHOLDER.CENTER_TITLE or ph.type == PP_PLACEHOLDER.VERTICAL_TITLE:
                            titleObj.append((str(shape.text_frame.text).lower()))
            listObj['slides'] = titleObj
            return listObj
        except PackageNotFoundError, e:
            return "PPTX not Found.", e

    def getKeywords(self):
        keyList = list()
        for key in self.getTitle()['slides']:
            keyList.append(k.get_keywords(key))
        return keyList
#         return keyObj
# if __name__ == '__main__':
#     p = ProcessSlide('slides/lecture2.pptx')
#     print p.getKeywords()
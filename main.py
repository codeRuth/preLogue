from pptx import Presentation
from pptx.enum.shapes import PP_PLACEHOLDER
from pptx.exc import PackageNotFoundError
import sys
import json


def get_title_object(input_presentation):
    main_obj = list()
    list_obj = dict()
    try :
        prs = Presentation(input_presentation)
        for index, slide in enumerate(prs.slides):
            for indexSlide, shape in enumerate(slide.shapes):
                if shape.is_placeholder:
                    phf = shape.placeholder_format
                    if phf.type == PP_PLACEHOLDER.TITLE or phf.type == PP_PLACEHOLDER.CENTER_TITLE or phf.type == PP_PLACEHOLDER.VERTICAL_TITLE:
                        main_obj.append(str(shape.text_frame.text))
        list_obj['slides'] = main_obj
        return json.dumps(list_obj)
    except PackageNotFoundError:
        print("PPTX not Found.")


def main():
    try :
        str(sys.argv[1])
        print get_title_object(str(sys.argv[1]))
    except IndexError:
        print("File Argument not given.")


if __name__ == "__main__":

    # if __name__ == '__main__':
    #     with open('stoplist.txt') as f:
    #         stopwords = f.read().splitlines()
    #
    #     queries = ["The mobile web is more important than mobile apps.",
    #                "As a #roadmapscholar, I highly recommend #startup bootcamp for #founders by @andrewsroadmaps : http://t.co/ZBISIMEBRH (http://t.co/VF5CojRWNF)",
    #                "RT @andrewsroadmaps: Proud of @naushadzaman &amp; @WasimKhal for winning the #IBMWatson hackathon! #roadmapscholars https://t.co/08sbAjKWKu."]
    #
    #     for query in queries:
    #         print 'query = ', query
    #         print 'get_keywords(query) = ', get_keywords(query)
    #         print "get_keywords(query, ['NP'])", get_keywords(query, ['NP'])
    #         print "extract_hashtag(query)", extract_hashtag(query)
    #         print "extract_users(query)", extract_users(query)
    #         print "extract_link(query)", extract_link(query)
    #         print ''

    main()
